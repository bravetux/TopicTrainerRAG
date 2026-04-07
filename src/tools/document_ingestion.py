"""Document ingestion pipeline: parse, chunk, embed, and index into ChromaDB."""
import os
import logging
import argparse
from pathlib import Path

import chromadb
from chromadb.config import Settings
from langchain_text_splitters import RecursiveCharacterTextSplitter

from src.config import (
    CHROMA_PERSIST_DIR, BUILTIN_TOPICS,
    CHUNK_SIZE, CHUNK_OVERLAP,
)
from src.tools.embedding_manager import embed_texts, get_embedding_config

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md"}


def get_supported_extensions() -> set:
    return SUPPORTED_EXTENSIONS


def parse_document(file_path: str) -> str:
    """Parse a document and return its full text content."""
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file format: {ext}. Supported: {SUPPORTED_EXTENSIONS}")

    if ext in (".txt", ".md"):
        return path.read_text(encoding="utf-8", errors="ignore")

    if ext == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)

    if ext == ".docx":
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)

    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            if texts:
                slides.append("\n".join(texts))
        return "\n\n".join(slides)

    if ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        sheets = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheets.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets)

    return ""


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping chunks."""
    if not text.strip():
        return []
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        length_function=len,
    )
    return splitter.split_text(text)



def get_chroma_client() -> chromadb.PersistentClient:
    """Return a persistent ChromaDB client."""
    return chromadb.PersistentClient(
        path=CHROMA_PERSIST_DIR,
        settings=Settings(anonymized_telemetry=False),
    )


def index_directory(
    directory: str,
    collection_name: str,
    chroma_client: chromadb.PersistentClient,
    reindex: bool = False,
    cancel_event=None,
    on_file_start=None,
    bedrock_client=None,  # kept for backwards-compatibility, ignored
) -> int:
    """Parse, chunk, embed, and store all documents in a directory. Returns chunk count.

    Args:
        cancel_event: Optional threading.Event. Indexing stops cleanly between files when set.
        on_file_start: Optional callable(file_name, file_index, total_files) called before each file.
    """
    collection = chroma_client.get_or_create_collection(
        name=collection_name,
        metadata={"hnsw:space": "cosine"},
    )

    if reindex:
        chroma_client.delete_collection(collection_name)
        collection = chroma_client.create_collection(
            name=collection_name,
            metadata={"hnsw:space": "cosine"},
        )

    total_chunks = 0
    doc_dir = Path(directory)
    files = [f for f in doc_dir.rglob("*") if f.suffix.lower() in SUPPORTED_EXTENSIONS]
    emb_cfg = get_embedding_config()

    for idx, file_path in enumerate(files):
        if cancel_event is not None and cancel_event.is_set():
            logger.info("Indexing cancelled after %d chunks.", total_chunks)
            break
        if on_file_start is not None:
            on_file_start(file_path.name, idx, len(files))
        logger.info("Parsing %s...", file_path.name)
        try:
            text = parse_document(str(file_path))
            if not text.strip():
                logger.warning("Empty document: %s", file_path.name)
                continue
            chunks = chunk_text(text)
            if not chunks:
                continue
            embeddings = embed_texts(chunks, emb_cfg)
            ids = [f"{file_path.stem}__chunk{i}" for i in range(len(chunks))]
            metadatas = [
                {
                    "source_file": file_path.name,
                    "collection": collection_name,
                    "chunk_index": i,
                }
                for i in range(len(chunks))
            ]
            collection.upsert(
                ids=ids,
                embeddings=embeddings,
                documents=chunks,
                metadatas=metadatas,
            )
            total_chunks += len(chunks)
            logger.info("  → %d chunks indexed", len(chunks))
        except Exception as exc:
            logger.error("Failed to index %s: %s", file_path.name, exc)

    return total_chunks


def _find_topic(topic_id: str) -> dict | None:
    """Find a topic by id — built-ins first, then custom registry."""
    for t in BUILTIN_TOPICS:
        if t["id"] == topic_id:
            return t
    try:
        from src.tools.kb_manager import _load_registry
        registry = _load_registry()
        for t in registry.get("custom", []):
            if t["id"] == topic_id:
                return t
    except Exception:
        pass
    return None


def index_technology(
    topic_id: str,
    reindex: bool = False,
    cancel_event=None,
    on_file_start=None,
) -> int:
    """Index documents for a single technology topic. Returns chunk count.

    Args:
        topic_id: Topic identifier (e.g. 'selenium', 'kubernetes').
        reindex: If True, drops the existing collection and rebuilds from scratch.
        cancel_event: Optional threading.Event. Stops cleanly between files when set.
        on_file_start: Optional callable(file_name, file_index, total_files).

    Returns:
        Number of chunks indexed. Returns 0 if topic not found or directory is empty.
    """
    topic = _find_topic(topic_id)
    if topic is None:
        logger.error("Unknown topic: %s", topic_id)
        return 0

    chroma = get_chroma_client()
    logger.info("Indexing topic '%s' (reindex=%s)...", topic_id, reindex)
    chunks = index_directory(
        directory=topic["doc_dir"],
        collection_name=topic["collection"],
        chroma_client=chroma,
        reindex=reindex,
        cancel_event=cancel_event,
        on_file_start=on_file_start,
    )
    logger.info("Topic '%s' indexed: %d chunks", topic_id, chunks)
    return chunks


def run_ingestion(reindex: bool = False) -> None:
    """Run the full ingestion pipeline for all built-in topics."""
    logger.info("Starting ingestion for all built-in topics (reindex=%s)...", reindex)
    total = 0
    for topic in BUILTIN_TOPICS:
        chunks = index_technology(topic["id"], reindex=reindex)
        total += chunks
    logger.info("Ingestion complete. Total chunks indexed: %d", total)


def check_status() -> None:
    """Print current ChromaDB collection stats for all built-in topics."""
    chroma = get_chroma_client()
    for topic in BUILTIN_TOPICS:
        try:
            col = chroma.get_collection(topic["collection"])
            print(f"{topic['id']}: {col.count()} chunks")
        except Exception:
            print(f"{topic['id']}: not indexed (run ingestion)")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--reindex", action="store_true", help="Delete and rebuild all collections")
    parser.add_argument("--status", action="store_true", help="Show current collection stats")
    args = parser.parse_args()

    if args.status:
        check_status()
    else:
        run_ingestion(reindex=args.reindex)
